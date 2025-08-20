import pandas as pd
import os
from flask import Flask, request, jsonify, render_template
from werkzeug.utils import secure_filename
import sys
import threading
import webbrowser
import re
from rapidfuzz import process, fuzz
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement

def get_base_dir():
    if getattr(sys, 'frozen', False):
        return sys._MEIPASS  # PyInstaller extracts to this temp dir
    return os.path.dirname(os.path.abspath(__file__))

app = Flask(__name__,
            template_folder=os.path.join(get_base_dir(), 'templates'))

BASE_DIR = get_base_dir()
print("Base Directory:", BASE_DIR)
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)


def remove_slide(prs, slide):
    # Helper to remove a slide from a Presentation (python-pptx workaround)
    slide_id = slide.slide_id
    slides = prs.slides._sldIdLst
    for sldId in slides:
        if int(sldId.attrib['id']) == slide_id:
            slides.remove(sldId)
            break



def apply_border(cell, edges = ["left", "right", "top", "bottom"], border_color="000000", border_width=1):
    if type(edges) is not list: edges = [edges]
    border_width = str(border_width*Pt(1))
    def SubElement(parent, tagname, **kwargs):
            element = OxmlElement(tagname)
            element.attrib.update(kwargs)
            parent.append(element)
            return element
    
    lines = [{"left": 'a:lnL',
              "right": 'a:lnR',
              "top": 'a:lnT',
              "bottom": 'a:lnB'}[_] for _ in edges]
    
    if cell.fill.type == MSO_FILL.SOLID: fill_color = cell.fill.fore_color.rgb
    cell.fill.background()
    
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for line in lines:
        
        # Remove duplicate tag if it exists
        tag = line.split(":")[-1]
        for e in tcPr.getchildren():
            if tag in str(e.tag): tcPr.remove(e)
        
        ln = SubElement(tcPr, line, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
        
    cell.fill.solid()
    if cell.fill.type == MSO_FILL.SOLID: cell.fill.fore_color.rgb = fill_color
    return(cell)

def duplicate_slide(prs: Presentation, slide_number: int,rows: int):
    """
    Duplicate a slide, copying only tables (structure and cell colors, no text).
    :param prs: Presentation object.
    :param slide_number: 1-based index of the slide to duplicate.
    :return: The newly created slide.
    """
    index = slide_number - 1
    if index < 0 or index >= len(prs.slides):
        raise IndexError(f"Slide {slide_number} does not exist.")

    source_slide = prs.slides[index]
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Remove all shapes from the new slide
    for shape in list(new_slide.shapes):
        sp = shape._element
        new_slide.shapes._spTree.remove(sp)

    # Add the table from the source slide
    for shape in source_slide.shapes:
        if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
            table = shape.table
            rows, cols = rows, len(table.columns)
            new_table_shape = new_slide.shapes.add_table(rows, cols, shape.left, shape.top, shape.width, shape.height)
            new_table = new_table_shape.table
            # Ensure header row height matches the source header row height
            new_table.rows[0].height = table.rows[0].height

            # Copy the first row (header) exactly
            from pptx.enum.text import MSO_VERTICAL_ANCHOR
            for c in range(cols):
                src_cell = table.cell(0, c)
                dest_cell = new_table.cell(0, c)
                dest_cell.text = src_cell.text
                # Set vertical alignment to middle
                dest_cell.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                # Copy font properties from the first run of the first paragraph if available
                if src_cell.text_frame.paragraphs and src_cell.text_frame.paragraphs[0].runs:
                    src_run = src_cell.text_frame.paragraphs[0].runs[0]
                    if dest_cell.text_frame.paragraphs and dest_cell.text_frame.paragraphs[0].runs:
                        dest_run = dest_cell.text_frame.paragraphs[0].runs[0]
                        dest_run.font.size = src_run.font.size
                        dest_run.font.bold = src_run.font.bold
                        dest_run.font.name = src_run.font.name
                        dest_run.font.italic = src_run.font.italic
                        # Only copy RGB color if available, otherwise skip to avoid AttributeError
                        if src_run.font.color and src_run.font.color.type is not None:
                            try:
                                rgb = src_run.font.color.rgb
                                if rgb is not None:
                                    dest_run.font.color.rgb = rgb
                            except AttributeError:
                                # src_run.font.color may be a scheme color or other type without .rgb
                                pass
                # Copy fill color
                try:
                    fill = src_cell.fill
                    new_fill = dest_cell.fill
                    if fill.type is not None and fill.fore_color.rgb is not None:
                        new_fill.solid()
                        new_fill.fore_color.rgb = fill.fore_color.rgb
                    else:
                        new_fill.solid()
                        new_fill.fore_color.rgb = RGBColor(255, 255, 255)
                    dest_cell = apply_border(dest_cell, edges=["left", "right", "top", "bottom"], border_color="000000", border_width=2)
                except Exception:
                    pass

            # Copy the rest of the table structure (clear text for data rows)
            for r in range(1, rows):
                new_table.rows[r].height = table.rows[1].height
                for c in range(cols):
                    new_table.columns[c].width = table.columns[c].width
                    src_cell = table.cell(1, c)
                    dest_cell = new_table.cell(r, c)
                    dest_cell.text = ""  # Clear text
                    
                    # Copy fill color
                    try:
                        fill = src_cell.fill
                        new_fill = dest_cell.fill
                        if fill.type is not None and fill.fore_color.rgb is not None:
                            new_fill.solid()
                            new_fill.fore_color.rgb = fill.fore_color.rgb
                        else:
                            new_fill.solid()
                            new_fill.fore_color.rgb = RGBColor(255, 255, 255)
                        dest_cell = apply_border(dest_cell, edges=["left", "right", "top", "bottom"], border_color="000000", border_width=2)
                        dest_cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    except Exception:
                        pass
            break  # Only one table per slide is handled

    # Add all other shapes (except the table) from the source slide to the new slide
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in source_slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
            # Copy text boxes
            if shape.has_text_frame:
                textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                textbox_tf = textbox.text_frame
                # Remove all paragraphs so the text frame is truly empty
                while len(textbox_tf.paragraphs) > 0:
                    p = textbox_tf.paragraphs[0]
                    textbox_tf._element.remove(p._element)
                # Now add paragraphs from the source, skipping empty ones
                for p in shape.text_frame.paragraphs:
                    if not p.text.strip():
                        continue
                    new_p = textbox_tf.add_paragraph()
                    new_p.text = p.text
                    new_p.alignment = p.alignment
                    for i, run in enumerate(p.runs):
                        if i == 0:
                            dest_run = new_p.runs[0]
                        else:
                            dest_run = new_p.add_run()
                        dest_run.text = run.text
                        if run.font:
                            dest_font = dest_run.font
                            dest_font.name = run.font.name
                            dest_font.size = run.font.size
                            dest_font.bold = run.font.bold
                            dest_font.italic = run.font.italic
                            if run.font.color and run.font.color.type is not None:
                                dest_font.color.rgb = run.font.color.rgb
            # Copy pictures
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                from io import BytesIO
                image_stream = BytesIO(image_bytes)
                new_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)
            # Copy auto shapes (basic shapes)
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
                # Optionally copy fill color
                try:
                    if shape.fill.type is not None and shape.fill.fore_color.rgb is not None:
                        new_shape.fill.solid()
                        new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
                except Exception:
                    pass
    return prs

def load_editable_presentation(source_path: str, day: str , Date: str) -> Presentation:
    """
    Copies the reference PPTX, opens it as a Presentation object, and replaces
    all instances of 'Month' with the specified month (0=January, 11=December).

    Returns the editable Presentation object so more slides can be added.
    """

    # Open the copied presentation
    prs = Presentation(source_path)

    # Replace 'Month' in all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace("Placeholder", Date)

    return prs  # You can now add slides, then save later
def _norm(s: str) -> str:
    s = str(s).lower().strip()
    s = re.sub(r"[^\w\s'-]", " ", s)
    s = s.replace("-", " ")
    s = re.sub(r"\s+", " ", s)
    return s

def _parse_name(s: str):
    s = _norm(s)
    parts = s.split()
    if not parts:
        return {"first":"", "middles":[], "last":"", "tokens":set()}
    prefixes = {"dr","mr","mrs","ms","miss"}
    suffixes = {"jr","sr","ii","iii","iv"}
    parts = [p for p in parts if p not in prefixes|suffixes]
    if not parts:
        return {"first":"", "middles":[], "last":"", "tokens":set()}
    first = parts[0]
    last = parts[-1] if len(parts) > 1 else ""
    middles = parts[1:-1] if len(parts) > 2 else []
    return {"first": first, "middles": middles, "last": last, "tokens": set(parts)}

def _first_match_score(a_first: str, b_first: str) -> float:
    if not a_first or not b_first:
        return 0.0
    a_first = a_first.strip().lower()
    b_first = b_first.strip().lower()
    # initial vs full name
    if a_first[0] == b_first[0] and (len(a_first) == 1 or len(b_first) == 1):
        return 95.0
    # fallback similarity
    return float(fuzz.ratio(a_first, b_first))

def _last_match_score(a_last: str, b_last: str) -> float:
    if not a_last or not b_last:
        return 0.0
    a = a_last.replace("'", "")
    b = b_last.replace("'", "")
    return float(max(fuzz.ratio(a, b), fuzz.token_sort_ratio(a, b)))

def name_similarity(a: str, b: str, *, score_cutoff: float = 0.0) -> float:
    pa, pb = _parse_name(a), _parse_name(b)
    last = _last_match_score(pa["last"], pb["last"])
    if last < 75:
        return 0.0
    first = _first_match_score(pa["first"], pb["first"])
    middle = 0.0
    if pa["middles"] and pb["middles"]:
        middle = float(fuzz.token_set_ratio(" ".join(pa["middles"]), " ".join(pb["middles"])))
    full = float(fuzz.token_set_ratio(_norm(a), _norm(b)))
    weighted = 0.7*last + 0.25*first + 0.05*middle
    score = max(weighted, full if last >= 85 else weighted)
    return score if score >= score_cutoff else 0.0

def match_providers(providers, org_providers, threshold=76):
    matched = []
    unmatched_map = {}  # providers_norm -> best unmatched

    matched_providers_set = set()  # to avoid re-matching same provider
    already_matched_org_providers_names = set()

    providers_map = { _norm(p): p for p in providers }
    org_providers_map = { _norm(p): p for p in org_providers }
    comp_norm = list(providers_map.keys())
    cap_norm = list(org_providers_map.keys())

    def last_token(s):
        parts = _norm(s).split()
        return parts[-1] if parts else ""

    # simple blocking
    cap_by_last, cap_by_initial = {}, {}
    for n in cap_norm:
        lt = last_token(n)
        cap_by_last.setdefault(lt, []).append(n)
        cap_by_initial.setdefault(lt[:1], []).append(n)

    for cr in comp_norm:
        original_cr = providers_map[cr]
        if original_cr in matched_providers_set:
            continue

        best_name = None
        best_score = -1.0
        matched_here = False

        for current_threshold in range(90, threshold - 1, -5):
            lt = last_token(cr)
            cand = set(cap_by_last.get(lt, [])) | set(cap_by_initial.get(lt[:1], []))
            if not cand:
                cand = set(cap_norm)
            cand = [c for c in cand if c not in already_matched_org_providers_names]
            if not cand:
                continue

            best = process.extractOne(cr, cand, scorer=name_similarity)
            if best is None:
                continue

            match_name, score, _ = best
            score = float(score)

            if score > best_score:
                best_score = score
                best_name = match_name

            if score >= current_threshold:
                matched.append({
                    'providers': original_cr,
                    'org_providers_name': org_providers_map[match_name],
                    'score': score
                })
                matched_providers_set.add(original_cr)
                already_matched_org_providers_names.add(match_name)
                matched_here = True
                break

        if not matched_here:
            unmatched_map[cr] = {
                'providers': original_cr,
                'org_providers_name': org_providers_map[best_name] if best_name and best_score > 0 else None,
                'score': best_score if best_score > 0 else 0.0
            }

    unmatched = sorted(unmatched_map.values(), key=lambda x: x['score'], reverse=True)
    return {"matched": matched, "unmatched": unmatched}

def preprocess_PR_FDR(PR_df, FDR_df):
        # Merge first 2 rows into one header row
    header1 = PR_df.iloc[0].fillna("")
    header2 = PR_df.iloc[1].fillna("")

    # Build combined header
    new_cols = []
    for h1, h2 in zip(header1, header2):
        if h1 and h2:
            new_cols.append(f"{h1}_{h2}".strip())
        elif h1:  # only top header
            new_cols.append(h1.strip())
        else:     # only sub header
            new_cols.append(h2.strip())
    PR_df = PR_df[3:]
    PR_df.columns = new_cols
    PR_df = PR_df.reset_index(drop=True)

    # Optional: clean names
    PR_df.columns = (
        PR_df.columns.str.replace(r"\s+", "_", regex=True)
                 .str.replace(r"[^\w]", "", regex=True)
                 .str.replace("__", "_")
    )

    # take a copy of the slice to avoid SettingWithCopy
    PR_df = PR_df.loc[:, ["Personnel","Earnings_Reg","Gross"]].copy()
    PR_df.dropna(inplace=True)
    PR_df['Personnel'] = PR_df['Personnel'].str.split("\n").str[0]
    for row in PR_df.itertuples():
        personnel = row.Personnel
        personnel = personnel.strip().split(",")
        personnel = personnel[1].strip() + " " + personnel[0].strip()
        PR_df.at[row.Index, 'Personnel'] = personnel
    def _to_numeric_series(s: pd.Series) -> pd.Series:
        cleaned = (
            s.astype(str).str.strip()
             .str.replace(r'[,\s$]', '', regex=True)
             .str.replace(r'^\((.*)\)$', r'-\1', regex=True)
        )
        return pd.to_numeric(cleaned, errors='coerce')

    # Ensure numeric types for payroll
    PR_df["Earnings_Reg"] = _to_numeric_series(PR_df["Earnings_Reg"])
    PR_df["Gross"] = _to_numeric_series(PR_df["Gross"])

    # FDR cleanup and numeric (copy slice first)
    FDR_df = FDR_df.loc[:, ["PROVIDER","NET CHARGES","NET RECEIPTS"]].copy()
    FDR_df["NET CHARGES"] = _to_numeric_series(FDR_df["NET CHARGES"])
    FDR_df["NET RECEIPTS"] = _to_numeric_series(FDR_df["NET RECEIPTS"])

    FDR_df = FDR_df[~((FDR_df["NET CHARGES"] == 0.0) & (FDR_df["NET RECEIPTS"] == 0.0))]
    FDR_df.drop_duplicates(subset=["PROVIDER"], inplace=True)
    PR_df.drop_duplicates(subset=["Personnel"], inplace=True)
    return PR_df, FDR_df


def build_metadata(charge_capture_df, company_roster_df, PR_df, FDR_df):
    CC_used_col = ["Provider", "CPT Codes", "Charge Status"]
    CR_used_col = ['Name', "Manager", 'State/Region']

    CC_filtered_df = charge_capture_df[CC_used_col]
    CR_filtered_df = company_roster_df[CR_used_col]


    CR_filtered_df = CR_filtered_df.sort_values(
        by='Name', key=lambda s: s.astype(str).str.lower(), kind='mergesort'
    )

    # Rename column safely
    CR_filtered_df = CR_filtered_df.rename(columns={'Name': 'Provider'})

    grouped_CC = CC_filtered_df.groupby('Provider')
    grouped_CR = CR_filtered_df.groupby('State/Region')
    cpt_pattern = re.compile(r'\b993\d{2}\b')  # Example pattern for 5-digit codes starting with 9930
    Margin_df = pd.DataFrame({
        'Clinician_Name': pd.Series(dtype='object'),
        'Manager': pd.Series(dtype='object'),
        'Gross_Consents': pd.Series(dtype='int'),
        'Gross_Encounters': pd.Series(dtype='int'),
        'Drafted_Encounters': pd.Series(dtype='int'), 
        'Collections': pd.Series(dtype='float'), 
        'Anticipated_Collections': pd.Series(dtype='float'),
        'Regular_Pay': pd.Series(dtype='float'),
        'Regular_Margin': pd.Series(dtype='float'),
        'Gross_Pay': pd.Series(dtype='float'),
        'Net_Margin': pd.Series(dtype='float'),
        'Anticipated_Net_Margin': pd.Series(dtype='float'),
        'region': pd.Series(dtype='object'),
    })
    for name, group in grouped_CC:
        group = group.copy()
        group['CPT Codes'] = group['CPT Codes'].astype(str).str.split(',')
        group = group.explode('CPT Codes')
        group['CPT Codes'] = group['CPT Codes'].str.strip()
        group['Charge Status'] = group['Charge Status'].str.strip().str.lower()
        # CCM Counts
        CCM_counts =group[group['CPT Codes']=="44444"]
        CCM_counts = CCM_counts['CPT Codes'].count()

        # Filter rows where 'CPT Codes' is in the target list and 'Charge Status' is 'draft'
        draft_counts = group[
            (group['CPT Codes'].str.match(cpt_pattern, na=False)) & 
            (group['Charge Status'] == 'draft')
        ]

        # Count the rows
        draft_counts = draft_counts['Charge Status'].count()
        # Gross Encounters    
        cpt_grouped = group[group['CPT Codes'].str.match(cpt_pattern, na=False)]
        gross_encounters = cpt_grouped['CPT Codes'].count()

        new_row = {'Clinician_Name': name,
                    'Manager': None,
                    'Gross_Consents': CCM_counts,
                    'Gross_Encounters': gross_encounters,
                    'Drafted_Encounters': draft_counts,
                    'Collections': 0.0,  # Placeholder for Collections
                    'Anticipated_Collections': 0.0,  # Placeholder for Anticipated Collections
                    'Regular_Pay': 0.0,  # Placeholder for Regular Pay
                    'Regular_Margin': 0.0,  # Placeholder for Regular Margin
                    'Gross_Pay': 0.0,  # Placeholder for Gross Pay
                    'Net_Margin': 0.0,  # Placeholder for Net Margin
                    'Anticipated_Net_Margin': 0.0,  # Placeholder for Anticipated Net Margin
                    "region": None,
                    }   

        Margin_df.loc[len(Margin_df)] = new_row
    
    
    PR_names = PR_df['Personnel'].tolist()
    FDR_names = FDR_df['PROVIDER'].tolist()
    Margin_df_names= Margin_df['Clinician_Name'].tolist()

    # Match providers from PR with Margin_df
    PR_result = match_providers(PR_names, Margin_df_names, threshold=85)
    matched_providers_PR = PR_result['matched']
    unmatched_providers_PR = PR_result['unmatched'] 
    print(f"Matched Providers from PR: {len(matched_providers_PR)}")
    print(f"Unmatched Providers from PR: {len(unmatched_providers_PR)}")
    # Update Margin_df with PR data
    for row in matched_providers_PR:
        clinician_name = row['providers']
        charge_capture_name = row['org_providers_name']
        score = row['score']

        reg_series = PR_df.loc[PR_df['Personnel'] == clinician_name, "Earnings_Reg"]
        reg_value = float(reg_series.iloc[0]) if not reg_series.empty and pd.notna(reg_series.iloc[0]) else np.nan
        Margin_df.loc[Margin_df['Clinician_Name'] == charge_capture_name, 'Regular_Pay'] = reg_value

        gross_series = PR_df.loc[PR_df['Personnel'] == clinician_name, "Gross"]
        gross_value = float(gross_series.iloc[0]) if not gross_series.empty and pd.notna(gross_series.iloc[0]) else np.nan
        Margin_df.loc[Margin_df['Clinician_Name'] == charge_capture_name, 'Gross_Pay'] = gross_value


    # Match providers from FDR with Margin_df
    FDR_result = match_providers(FDR_names, Margin_df_names, threshold=85)
    matched_providers_FDR = FDR_result['matched']
    unmatched_providers_FDR = FDR_result['unmatched']
    print(f"Matched Providers from FDR: {len(matched_providers_FDR)}")
    print(f"Unmatched Providers from FDR: {len(unmatched_providers_FDR)}")

    # Update Margin_df with FDR data
    for row in matched_providers_FDR:
        clinician_name = row['providers']
        charge_capture_name = row['org_providers_name']
        score = row['score']

        # Assign scalars, not Series
        charges_series = FDR_df.loc[FDR_df['PROVIDER'] == clinician_name, "NET CHARGES"]
        charges_value = float(charges_series.iloc[0]) if not charges_series.empty and pd.notna(charges_series.iloc[0]) else np.nan
        Margin_df.loc[Margin_df['Clinician_Name'] == charge_capture_name, 'Anticipated_Collections'] = charges_value

        receipts_series = FDR_df.loc[FDR_df['PROVIDER'] == clinician_name, "NET RECEIPTS"]
        receipts_value = float(receipts_series.iloc[0]) if not receipts_series.empty and pd.notna(receipts_series.iloc[0]) else np.nan
        Margin_df.loc[Margin_df['Clinician_Name'] == charge_capture_name, 'Collections'] = receipts_value
        # Regular_Margin "(regular pay column / collections)*100"
        if not pd.isna(receipts_value) and receipts_value != 0:
            regular_pay = Margin_df.loc[Margin_df['Clinician_Name'] == charge_capture_name, 'Regular_Pay'].values[0]
            if not pd.isna(regular_pay):
                regular_margin = (receipts_value / regular_pay ) * 100 if regular_pay != 0 else 0
                Margin_df.loc[Margin_df['Clinician_Name'] == charge_capture_name, 'Regular_Margin'] = round(regular_margin,2)
                
        # Net_Margin "(gross pay column / collections)*100"
        if not pd.isna(charges_value) and charges_value != 0:
            gross_pay = Margin_df.loc[Margin_df['Clinician_Name'] == charge_capture_name, 'Gross_Pay'].values[0]
            if not pd.isna(gross_pay):
                Anticipated_Net_Margin = (charges_value/ gross_pay ) * 100 if gross_pay != 0 else 0
                net_margin = (receipts_value / gross_pay) * 100 if gross_pay != 0 else 0
                Margin_df.loc[Margin_df['Clinician_Name'] == charge_capture_name, 'Anticipated_Net_Margin'] = round(Anticipated_Net_Margin, 2)
                Margin_df.loc[Margin_df['Clinician_Name'] == charge_capture_name, 'Net_Margin'] = round(net_margin, 2)

    

    unmatched_providers = []
    matched_providers = []
    Regional_Dashboard=pd.DataFrame(columns=["Region","RCS","RDO","Gross Encounters","Gross Consents","Gross Drafts","Regular_margin","Gross_margin","Anticipated_Net_Margin"])

    for name, group in grouped_CR:
        manager_group=group.groupby('Manager')
        for manager_name, manager_group in manager_group:
            Region_list=Margin_df['Clinician_Name']
            manager_group_list = manager_group['Provider']
            result=match_providers(manager_group_list, Region_list, threshold=85)
            unmatched_providers.extend(result['unmatched'])
            matched_providers.extend(result['matched'])
            matched_providers_CC = [item['org_providers_name'] for item in result['matched']]
            Region_Gross_Encounters = Margin_df[Margin_df['Clinician_Name'].isin([item['org_providers_name'] for item in result['matched']])]['Gross_Encounters'].sum()
            Region_Gross_Consents = Margin_df[Margin_df['Clinician_Name'].isin([item['org_providers_name'] for item in result['matched']])]['Gross_Consents'].sum()
            Region_Gross_Drafts = Margin_df[Margin_df['Clinician_Name'].isin([item['org_providers_name'] for item in result['matched']])]['Drafted_Encounters'].sum()
            Region_Regular_margin = Margin_df[Margin_df['Clinician_Name'].isin([item['org_providers_name'] for item in result['matched']])]['Regular_Margin'].mean()
            Region_Gross_margin = Margin_df[Margin_df['Clinician_Name'].isin([item['org_providers_name'] for item in result['matched']])]['Net_Margin'].mean()
            Region_Anticipated_Net_Margin = Margin_df[Margin_df['Clinician_Name'].isin([item['org_providers_name'] for item in result['matched']])]['Anticipated_Net_Margin'].mean()  
            # Create a new row for the Regional Dashboard
            if not matched_providers_CC:
                continue  # Skip if no matched providers  
            new_row = {
                "Region": name,
                "RCS": manager_name,
                "RDO": matched_providers_CC,
                "Gross Encounters": Region_Gross_Encounters,
                "Gross Consents": Region_Gross_Consents,
                "Gross Drafts": Region_Gross_Drafts,
                "Regular_margin": Region_Regular_margin,  # Placeholder for Net Margin
                "Gross_margin": Region_Gross_margin,  # Placeholder for Gross Margin
                "Anticipated_Net_Margin": Region_Anticipated_Net_Margin,  # Placeholder for Anticipated Net Margin
            }
            Regional_Dashboard.loc[len(Regional_Dashboard)] = new_row
    Regional_Dashboard.columns = Regional_Dashboard.columns.str.replace(' ', '_')
    Margin_df.columns = Margin_df.columns.str.replace(' ', '_')
    for row in Regional_Dashboard.itertuples():
        region = row.Region
        Clinician_list = row.RDO
        manager= row.RCS    
        # Update the 'region' column for all clinicians in the list
        Margin_df.loc[Margin_df['Clinician_Name'].isin(Clinician_list), 'region'] = region
        Margin_df.loc[Margin_df['Clinician_Name'].isin(Clinician_list), 'Manager'] = manager

        # Sort alphabetically (case-insensitive) as requested
    Margin_df = Margin_df.sort_values(
        by='Clinician_Name', key=lambda s: s.astype(str).str.lower(), kind='mergesort'
    )
    # Drop rows with NaN values in the 'region' column
    Margin_df.dropna(subset=['region'], inplace=True)
    
    grouped_CC_ByRegion = Margin_df.groupby(['region', 'Manager'])
    return Regional_Dashboard, unmatched_providers, matched_providers, grouped_CC_ByRegion,Margin_df

def generate_pbj_presentation(prs:Presentation,Regional_Dashboard: pd.DataFrame,gouped_CC_ByRegion, day:str, date:str,Add_region=True):
    
    # Split the corporate summary table into slides of 10 facilities each
    Region_per_slide = 5
    num_Region = len(Regional_Dashboard)
    # Totals across all regions
    enouncter_sum = 0
    draft_sum = 0
    consent_sum = 0
    _reg_margin_sum = 0.0
    _reg_margin_cnt = 0
    _gross_margin_sum = 0.0
    _gross_margin_cnt = 0
    _anticipated_net_margin_sum = 0.0
    _anticipated_net_margin_cnt = 0

    for start_idx in range(0, num_Region, Region_per_slide):
        end_idx = min(start_idx + Region_per_slide, num_Region)
        is_last_chunk = end_idx == num_Region
        # +1 for header, +1 extra for totals row on the last slide
        rows = (end_idx - start_idx) + 1 + (1 if is_last_chunk else 0)
        prs = duplicate_slide(prs, 2, rows)
        region_dashborad_slide = prs.slides[-1]
        for shape in region_dashborad_slide.shapes:
            if shape.has_table:
                table = shape.table
                fontSize = table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size
                from pptx.enum.text import MSO_VERTICAL_ANCHOR
                for idx, row in enumerate(Regional_Dashboard.iloc[start_idx:end_idx].itertuples()):
                    # Use safe attribute access since namedtuple fields can vary
                    region = getattr(row, 'Region', '')
                    manager = getattr(row, 'RCS', '')
                    print(f"Processing Region: {region}, Manager: {manager}")
                    providers_list = getattr(row, 'RDO', [])
                    # Ensure providers_list is iterable
                    if providers_list is None:
                        providers_list = []
                    if isinstance(providers_list, (str, int, float)):
                        providers_list = [providers_list]
                    gross_encounters = getattr(row, 'Gross_Encounters', '')
                    gross_consents = getattr(row, 'Gross_Consents', '')
                    gross_drafts = getattr(row, 'Gross_Drafts', '')
                    # Capture raw values for mean calculations (before rounding/empty handling)
                    _reg_val = getattr(row, 'Regular_margin', np.nan)
                    _gross_val = getattr(row, 'Gross_margin', np.nan)
                    _anticip_val = getattr(row, 'Anticipated_Net_Margin', np.nan)
                    Regular_margin = round(_reg_val, 2) if not pd.isna(_reg_val) else ''
                    Gross_margin = round(_gross_val, 2) if not pd.isna(_gross_val) else ''
                    Anticipated_Net_Margin = round(_anticip_val, 2) if not pd.isna(_anticip_val) else ''
                    row_cells = table.rows[idx+1].cells
                    row_cells[0].text = str(region)
                    row_cells[1].text = str(manager)
                    # Split providers_list into two columns, pad each provider name to fixed width, and set monospace font
                    monospace_font = "Consolas"
                    pad_width = 25
                    names = [str(name) if len(str(name)) > pad_width else str(name).ljust(pad_width) for name in providers_list]
                    col1 = []
                    col2 = []
                    for i, name in enumerate(names):
                        if i % 2 == 0:
                            col1.append(name)
                        else:
                            col2.append(name)
                    # Make both columns the same length
                    max_len = max(len(col1), len(col2))
                    while len(col1) < max_len:
                        col1.append(' ' * pad_width)
                    while len(col2) < max_len:
                        col2.append('')
                    lines = [col1[i] + col2[i] for i in range(max_len)]
                    combined_text = '\n'.join(lines)
                    row_cells[2].text = combined_text
                    # Set monospace font for all runs in the cell
                    for para in row_cells[2].text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = monospace_font
                    row_cells[3].text = str(gross_encounters)
                    row_cells[4].text = str(gross_consents)
                    row_cells[5].text = str(gross_drafts)
                    row_cells[6].text = str(Regular_margin)+" %"
                    row_cells[7].text = str(Gross_margin)+" %"
                    row_cells[8].text = str(Anticipated_Net_Margin)+" %"
                    # Set vertical alignment to middle for all cells and all paragraphs in each cell
                    for cell in row_cells:
                        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.size = fontSize

                    # Accumulate totals (safe numeric conversions)
                    try:
                        if gross_encounters is not None and not pd.isna(gross_encounters):
                            enouncter_sum += float(gross_encounters)
                    except Exception:
                        pass
                    try:
                        if gross_drafts is not None and not pd.isna(gross_drafts):
                            draft_sum += float(gross_drafts)
                    except Exception:
                        pass
                    try:
                        if gross_consents is not None and not pd.isna(gross_consents):
                            consent_sum += float(gross_consents)
                    except Exception:
                        pass
                    if _reg_val is not None and not pd.isna(_reg_val):
                        _reg_margin_sum += float(_reg_val)
                        _reg_margin_cnt += 1
                    if _gross_val is not None and not pd.isna(_gross_val):
                        _gross_margin_sum += float(_gross_val)
                        _gross_margin_cnt += 1
                    if _anticip_val is not None and not pd.isna(_anticip_val):
                        _anticipated_net_margin_sum += float(_anticip_val)
                        _anticipated_net_margin_cnt += 1

                # If this is the last chunk, write totals to the last row of this slide
                if is_last_chunk:
                    total_row_idx = (end_idx - start_idx) + 1  # header + data rows, totals is next
                    total_cells = table.rows[total_row_idx].cells
                    total_cells[0].text = "Practice Total"
                    total_cells[1].text = ""
                    total_cells[2].text = ""
                    total_cells[3].text = str(int(enouncter_sum))
                    total_cells[4].text = str(int(consent_sum))
                    total_cells[5].text = str(int(draft_sum))
                    reg_mean = round(_reg_margin_sum / _reg_margin_cnt, 2) if _reg_margin_cnt else ''
                    gross_mean = round(_gross_margin_sum / _gross_margin_cnt, 2) if _gross_margin_cnt else ''
                    anticip_mean = round(_anticipated_net_margin_sum / _anticipated_net_margin_cnt, 2) if _anticipated_net_margin_cnt else ''
                    total_cells[6].text = str(reg_mean)+ " %"
                    total_cells[7].text = str(gross_mean)+ " %"
                    total_cells[8].text = str(anticip_mean)+ " %"
                    # Apply same vertical alignment and font size
                    for cell in total_cells:
                        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.size = fontSize+20
                                run.font.bold = True  # Make totals bold
    for region_name, region_data in gouped_CC_ByRegion:
        # --- SPLIT PROVIDERS ACROSS SLIDES BASED ON MAX PROVIDERS ---
        max_providers_per_slide = 7
        providers_df = region_data.reset_index(drop=True)
        total_providers = len(providers_df)
        print(f"Region: {region_name}, Total Providers: {total_providers}")

        # Create as many slides as needed to fit all providers
        for start in range(0, total_providers, max_providers_per_slide):
            end = min(start + max_providers_per_slide, total_providers)
            rows = (end - start) + 1  # +1 for header
            prs = duplicate_slide(prs, 3, rows)
            Region_Scope_slide = prs.slides[-1]

            for shape in Region_Scope_slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = run.text.replace("Region", f"{region_name[0]} - Manager: {region_name[1]}")
                if shape.has_table:
                    table = shape.table
                    fontSize = table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size
                    font = table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.name
                    from pptx.enum.text import MSO_VERTICAL_ANCHOR

                    # Fill only the current chunk of providers
                    for idx, row in enumerate(providers_df.iloc[start:end].itertuples(index=False)):
                        # Access expected Margin_df columns safely
                        clinician = getattr(row, 'Clinician_Name', None) or getattr(row, 'Clinician', '')
                        gross_encounters = getattr(row, 'Gross_Encounters', '')
                        gross_consents = getattr(row, 'Gross_Consents', '')
                        gross_drafts = getattr(row, 'Drafted_Encounters', '')

                        row_cells = table.rows[idx + 1].cells
                        row_cells[0].text = str(clinician)
                        row_cells[1].text = str(gross_encounters)
                        row_cells[2].text = str(gross_drafts)
                        row_cells[3].text = str(gross_consents)
                        row_cells[4].text = str(getattr(row, 'Regular_Margin', ''))+ " %"
                        row_cells[5].text = str(getattr(row, 'Net_Margin', ''))+ " %"
                        row_cells[6].text = str(getattr(row, 'Anticipated_Net_Margin', ''))+ " %"

                        # Set vertical alignment to middle for all cells and font size
                        for cell in row_cells:
                            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    run.font.size = fontSize
    # Safely remove template slides if they exist (ensure at least 2 slides)
    try:
        if len(prs.slides) > 1:
            remove_slide(prs, prs.slides[1])
        if len(prs.slides) > 1:
            remove_slide(prs, prs.slides[1])
    except Exception:
        # Non-fatal: if removal fails, continue
        pass
def save_workbook(workbook, filename):
    """Save the workbook (a pandas.DataFrame) to an Excel file with reasonable formatting.

    Parameters:
        workbook: pandas.DataFrame
        filename: str - output path
    """
    output_path = filename

    # Make a safe copy and stringify any list-like cells
    _to_str = lambda v: ", ".join(map(str, v)) if isinstance(v, (list, tuple)) else v
    _margin_out = workbook.copy()
    for col in _margin_out.select_dtypes(include=["object"]).columns:
        _margin_out[col] = _margin_out[col].map(_to_str)

    with pd.ExcelWriter(output_path, engine="xlsxwriter") as writer:
        _margin_out.to_excel(writer, sheet_name="Margin", index=False)
        wb = writer.book
        ws = writer.sheets["Margin"]

        header_fmt = wb.add_format({"bold": True, "text_wrap": True, "valign": "top", "fg_color": "#DDEBF7", "border": 1})
        wrap_fmt = wb.add_format({"text_wrap": True, "valign": "top"})
        int_fmt = wb.add_format({"num_format": "#,##0"})
        float_fmt = wb.add_format({"num_format": "#,##0.00"})

        # Rewrite headers with format
        for j, col in enumerate(_margin_out.columns):
            ws.write(0, j, col, header_fmt)

        # Autosize columns and set formats
        for j, col in enumerate(_margin_out.columns):
            ser = _margin_out[col]
            try:
                est_len = int(ser.astype(str).str.len().quantile(0.95)) if len(ser) else 0
            except Exception:
                est_len = int(ser.astype(str).str.len().max() if len(ser) else 0)
            max_len = max(len(str(col)), est_len) + 2
            width = min(max_len, 60)

            if pd.api.types.is_integer_dtype(ser):
                fmt = int_fmt
            elif pd.api.types.is_float_dtype(ser):
                fmt = float_fmt
            else:
                fmt = wrap_fmt

            ws.set_column(j, j, width, fmt)

        # Freeze header and add autofilter
        ws.freeze_panes(1, 0)
        ws.autofilter(0, 0, len(_margin_out), len(_margin_out.columns) - 1)
@app.route('/')
def index():
    return render_template('Margin_report.html')


@app.route('/list_pr_sheets', methods=['POST'])
def list_pr_sheets():
    """Return list of sheet names for an uploaded PR Excel file.

    Accepts multipart form with file field 'PRFile'. Returns JSON {sheets: [...]}.
    """
    if 'PRFile' not in request.files:
        return jsonify({'error': 'No PRFile provided'}), 400
    pr_file = request.files['PRFile']
    filename = secure_filename(pr_file.filename or '')
    ext = os.path.splitext(filename)[1].lower()
    if ext not in ('.xls', '.xlsx'):
        # still try to read if extension missing, but otherwise return error
        return jsonify({'error': 'Invalid PR file type; expected .xls or .xlsx'}), 400
    try:
        # pandas can read from file-like objects
        xls = pd.ExcelFile(pr_file)
        sheets = xls.sheet_names
        return jsonify({'sheets': sheets})
    except Exception as e:
        # fallback: save temp and try again
        try:
            tmp_path = os.path.join(app.config['UPLOAD_FOLDER'], f"tmp_pr_{threading.get_ident()}{ext}")
            pr_file.seek(0)
            pr_file.save(tmp_path)
            xls = pd.ExcelFile(tmp_path)
            sheets = xls.sheet_names
            try:
                os.remove(tmp_path)
            except Exception:
                pass
            return jsonify({'sheets': sheets})
        except Exception as e2:
            return jsonify({'error': str(e2)}), 500
# Flask API endpoint to receive file and CPT dict
@app.route('/upload_data', methods=['POST'])
def upload_data():
    if 'Roster' not in request.files or 'captureFile' not in request.files:
        return jsonify({'error': 'Missing files'}), 400
    Roster_file = request.files['Roster']
    capture_file = request.files['captureFile']
    PR_file = request.files.get('PRFile', None)
    FDR_file = request.files.get('FDRFile', None)
    PR_sheet_name = request.form.get('PR_sheet_name', 'Payroll Register')
    capture_file_filename = secure_filename(capture_file.filename)
    Roster_file_filename = secure_filename(Roster_file.filename)
    PR_file_filename = secure_filename(PR_file.filename) if PR_file else None
    FDR_file_filename = secure_filename(FDR_file.filename) if FDR_file else None
    captuer_ext = os.path.splitext(capture_file_filename)[1].lower()
    roster_ext = os.path.splitext(Roster_file_filename)[1].lower()
    PR_ext = os.path.splitext(PR_file_filename)[1].lower() if PR_file_filename else None
    FDR_ext = os.path.splitext(FDR_file_filename)[1].lower() if FDR_file_filename else None
    if captuer_ext not in ['.csv', '.xlsx'] or roster_ext not in ['.csv', '.xlsx']:
        return jsonify({'error': 'Invalid file type'}), 400

    # Save files to a temp location
    # Save file to a temp location
    capture_temp_path = os.path.join(app.config['UPLOAD_FOLDER'], capture_file_filename)
    roster_temp_path = os.path.join(app.config['UPLOAD_FOLDER'], Roster_file_filename)
    PR_temp_path = os.path.join(app.config['UPLOAD_FOLDER'], PR_file_filename) if PR_file_filename else None
    FDR_temp_path = os.path.join(app.config['UPLOAD_FOLDER'], FDR_file_filename) if FDR_file_filename else None
    capture_file.save(capture_temp_path)
    Roster_file.save(roster_temp_path)
    PR_file.save(PR_temp_path) if PR_file_filename else None
    FDR_file.save(FDR_temp_path) if FDR_file_filename else None
    # Load DataFrame

    # Read charge capture file
    if capture_file_filename.endswith('.csv'):
        charge_capture_df = pd.read_csv(capture_temp_path)
    else:
        charge_capture_df = pd.read_excel(capture_temp_path)

    # Read roster file
    if Roster_file_filename.endswith('.csv'):
        company_roaster = pd.read_csv(roster_temp_path)
    else:
        company_roaster = pd.read_excel(roster_temp_path, skiprows=2)

    # Read PR file if provided
    if PR_file_filename:
        if PR_ext.endswith('.xls'):
            PR_df = pd.read_excel(PR_temp_path,sheet_name=PR_sheet_name,header=None)
        else:
            return jsonify({'error': 'Invalid PR file type plase upload with xls extension'}), 400
    if FDR_file_filename:
        if FDR_ext.endswith('.xlsx'):
            FDR_df = pd.read_excel(FDR_temp_path)
        else:
            return jsonify({'error': 'Invalid FDR file type plase upload with xlsx extension'}), 400
    
    PR_df, FDR_df = preprocess_PR_FDR(PR_df, FDR_df)
    # Get month and report_type from form

    day = request.form.get('day')
    date= request.form.get('date')
   
    include_regions = request.form.get('include_regions',True)

    Regional_Dashboard, unmatched_providers, matched_providers, grouped_CC_ByRegion,Margin_df = build_metadata(charge_capture_df, company_roaster,PR_df, FDR_df)
    prs = load_editable_presentation(os.path.join(BASE_DIR, 'static', 'Margin Report.pptx'), day= day, Date=date)
    generate_pbj_presentation(prs, Regional_Dashboard, grouped_CC_ByRegion, day, date, Add_region=include_regions)
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], 'Margin_report.pptx')
    workBook_path = os.path.join(app.config['UPLOAD_FOLDER'], 'Work_book.xlsx')
    # Save Margin_df to Excel
    save_workbook(Margin_df, workBook_path)
    # Save the presentation
    prs.save(output_path)
    # Optionally, return download link or status
    # Return the relative path for the frontend to use in download URL
    return jsonify({
        'success': True,
        'pptx_path': f"uploads/Margin_Report.pptx",
        'matched_providers': matched_providers,
        'unmatched_providers': unmatched_providers
    })


# from flask import send_from_directory

# # Serve files from the uploads directory
# @app.route('/uploads/<path:filename>')
# def download_file(filename):
#     return send_from_directory(app.config['UPLOAD_FOLDER'], filename, as_attachment=True)

# # Function to open the default web browser
# import os
# def open_browser():
#     webbrowser.open_new("http://localhost:8000")

# if __name__ == '__main__':
#     # Only open browser in local development
#     if os.environ.get('RENDER') is None:
#         threading.Timer(1, open_browser).start()  # Open the browser after 1 second
#     port = int(os.environ.get('PORT', 8000))
#     app.run(host='0.0.0.0', port=port,debug=True)