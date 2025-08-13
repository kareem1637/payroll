
import pandas as pd
import os
from flask import Flask, request, jsonify, render_template
from werkzeug.utils import secure_filename
import sys
import threading
import webbrowser
import math
from collections import defaultdict
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import  Pt
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement

def get_base_dir():
    if getattr(sys, 'frozen', False):
        return sys._MEIPASS  # PyInstaller extracts to this temp dir
    return os.path.dirname(os.path.abspath(__file__))

app = Flask(__name__,
            template_folder=os.path.join(get_base_dir(), 'templates'))

BASE_DIR = get_base_dir()
print("Base Directory:", BASE_DIR)
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)


def remove_slide(prs, slide):
    # Helper to remove a slide from a Presentation (python-pptx workaround)
    slide_id = slide.slide_id
    slides = prs.slides._sldIdLst
    for sldId in slides:
        if int(sldId.attrib['id']) == slide_id:
            slides.remove(sldId)
            break

def load_charge_capture_data(filepath):
    return pd.read_excel(filepath)

def group_facilities(charge_capture_df):
    sorted_fac = sorted(charge_capture_df['Facility'].unique())
    stopwords = {"of", "the", "at", "and", "care", "center", "rehab", "nursing", "health", "llc", "snf", "m", "by", "for", "on", "in", "living", "villa", "home"}
    grouped_facilities = defaultdict(list)
    for fac in sorted_fac:
        key = fac.split()[0]
        if key.lower() in stopwords:
            key = fac.split()[1]
        grouped_facilities[key].append(fac)
    corporate_groups = {k: v for k, v in grouped_facilities.items() if len(v) > 1}
    single_facilities = {k: v for k, v in grouped_facilities.items() if len(v) == 1}
    keys_to_remove = []
    for key, facilities in single_facilities.items():
        for corp_key in corporate_groups:
            if corp_key.lower() in single_facilities[key][0].lower():
                corporate_groups[corp_key].extend(facilities)
                keys_to_remove.append(key)
                break
    for key in keys_to_remove:
        if key in single_facilities:
            del single_facilities[key]
    return corporate_groups, single_facilities

def build_metadata(charge_capture_df, corporate_groups, single_facilities):
    cpt_time_dict = {
        "99304": 60,
        "99305": 75,
        "99306": 90,
        "99307": 15,
        "99308": 25,
        "99309": 35,
        "99310": 45
    }
    Facilites_groups = charge_capture_df.groupby('Facility', as_index=False)
    metadata = {}
    for key, facilities in corporate_groups.items():
        corp_data = pd.DataFrame()
        facility_summary = []
        facility_data_dict = {}
        for facility in facilities:
            group = Facilites_groups.get_group(facility).copy()
            group.loc[:, 'CPT Codes'] = group['CPT Codes'].astype(str).str.split(',')
            group = group.explode('CPT Codes')
            group['CPT Codes'] = group['CPT Codes'].str.strip()
            Facility_Level = group.groupby(['Provider', 'CPT Codes', 'DOS']).size().reset_index(name='Count')
            Facility_Level['Minutes'] = Facility_Level['CPT Codes'].astype(str).map(cpt_time_dict) * Facility_Level['Count']
            # Remove rows where Minutes == 0 or is NaN
            Facility_Level = Facility_Level[(Facility_Level['Minutes'] != 0) & (Facility_Level['Minutes'].notna())]
            if Facility_Level.empty:
                continue  # Skip this facility if no valid entries remain
            Facility_Level_result = Facility_Level.groupby(['Provider', 'DOS'], as_index=False)['Minutes'].sum()
            # Round up each clinician's hours to the nearest 0.5
            def round_up_half(x):
                return math.ceil(x * 2) / 2
            Facility_Level_result['Hours'] = Facility_Level_result['Minutes'].apply(lambda x: round_up_half(x / 60))
            # Group by Provider for this facility and sum hours (sum after rounding up each row)
            provider_hours = Facility_Level_result.groupby('Provider', as_index=False)['Hours'].sum()
            provider_hours['Facility'] = facility
            facility_summary.append(provider_hours)
            facility_data_dict[facility] = Facility_Level_result
        # Only add to metadata if there is at least one valid facility
        if not facility_summary:
            continue
        metadata[key] = {}
        metadata[key]['Facility_Data'] = facility_data_dict
        summary_df = pd.concat(facility_summary, ignore_index=True)
        # Group by Facility, aggregate clinicians and sum hours
        summary_grouped = summary_df.groupby('Facility').agg({
            'Provider': lambda x: '\n'.join(x),
            'Hours': 'sum'
        }).reset_index()
        summary_grouped = summary_grouped.rename(columns={'Provider': 'Clinician'})
        metadata[key]['Corp_data'] = summary_grouped
    for key, facilities in single_facilities.items():
        facility_name = facilities[0]  # Single facility case
        facility_name = facility_name.strip().split("and")[0]  # Ensure no leading/trailing spaces
        corp_data = pd.DataFrame()
        facility_summary = []
        facility_data_dict = {}
        for facility in facilities:
            group = Facilites_groups.get_group(facility).copy()
            group.loc[:, 'CPT Codes'] = group['CPT Codes'].astype(str).str.split(',')
            group = group.explode('CPT Codes')
            group['CPT Codes'] = group['CPT Codes'].str.strip()
            Facility_Level = group.groupby(['Provider', 'CPT Codes', 'DOS']).size().reset_index(name='Count')
            Facility_Level['Minutes'] = Facility_Level['CPT Codes'].astype(str).map(cpt_time_dict) * Facility_Level['Count']
            # Remove rows where Minutes == 0 or is NaN
            Facility_Level = Facility_Level[(Facility_Level['Minutes'] != 0) & (Facility_Level['Minutes'].notna())]
            if Facility_Level.empty:
                continue  # Skip this facility if no valid entries remain
            Facility_Level_result = Facility_Level.groupby(['Provider', 'DOS'], as_index=False)['Minutes'].sum()
            # Round up each clinician's hours to the nearest 0.5
            def round_up_half(x):
                return math.ceil(x * 2) / 2
            Facility_Level_result['Hours'] = Facility_Level_result['Minutes'].apply(lambda x: round_up_half(x / 60))
            # Group by Provider for this facility and sum hours (sum after rounding up each row)
            provider_hours = Facility_Level_result.groupby('Provider', as_index=False)['Hours'].sum()
            provider_hours['Facility'] = facility
            facility_summary.append(provider_hours)
            facility_data_dict[facility] = Facility_Level_result
        # Only add to metadata if there is at least one valid facility
        if not facility_summary:
            continue
        metadata[facility_name] = {}
        metadata[facility_name]['Facility_Data'] = facility_data_dict
        summary_df = pd.concat(facility_summary, ignore_index=True)
        # Group by Facility, aggregate clinicians and sum hours
        summary_grouped = summary_df.groupby('Facility').agg({
            'Provider': lambda x: '\n'.join(x),
            'Hours': 'sum'
        }).reset_index()
        summary_grouped = summary_grouped.rename(columns={'Provider': 'Clinician'})
        metadata[facility_name]['Corp_data'] = summary_grouped
    return metadata




def apply_border(cell, edges = ["left", "right", "top", "bottom"], border_color="000000", border_width=1):
    if type(edges) is not list: edges = [edges]
    border_width = str(border_width*Pt(1))
    def SubElement(parent, tagname, **kwargs):
            element = OxmlElement(tagname)
            element.attrib.update(kwargs)
            parent.append(element)
            return element
    
    lines = [{"left": 'a:lnL',
              "right": 'a:lnR',
              "top": 'a:lnT',
              "bottom": 'a:lnB'}[_] for _ in edges]
    
    if cell.fill.type == MSO_FILL.SOLID: fill_color = cell.fill.fore_color.rgb
    cell.fill.background()
    
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for line in lines:
        
        # Remove duplicate tag if it exists
        tag = line.split(":")[-1]
        for e in tcPr.getchildren():
            if tag in str(e.tag): tcPr.remove(e)
        
        ln = SubElement(tcPr, line, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
        
    cell.fill.solid()
    if cell.fill.type == MSO_FILL.SOLID: cell.fill.fore_color.rgb = fill_color
    return(cell)

def duplicate_slide(prs: Presentation, slide_number: int,rows: int):
    """
    Duplicate a slide, copying only tables (structure and cell colors, no text).
    :param prs: Presentation object.
    :param slide_number: 1-based index of the slide to duplicate.
    :return: The newly created slide.
    """
    index = slide_number - 1
    if index < 0 or index >= len(prs.slides):
        raise IndexError(f"Slide {slide_number} does not exist.")

    source_slide = prs.slides[index]
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Remove all shapes from the new slide
    for shape in list(new_slide.shapes):
        sp = shape._element
        new_slide.shapes._spTree.remove(sp)

    # Add the table from the source slide
    for shape in source_slide.shapes:
        if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
            table = shape.table
            rows, cols = rows, len(table.columns)
            new_table_shape = new_slide.shapes.add_table(rows, cols, shape.left, shape.top, shape.width, shape.height)
            new_table = new_table_shape.table

            # Copy the first row (header) exactly
            from pptx.enum.text import MSO_VERTICAL_ANCHOR
            for c in range(cols):
                src_cell = table.cell(0, c)
                dest_cell = new_table.cell(0, c)
                dest_cell.text = src_cell.text
                # Set vertical alignment to middle
                dest_cell.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                # Copy font properties from the first run of the first paragraph if available
                if src_cell.text_frame.paragraphs and src_cell.text_frame.paragraphs[0].runs:
                    src_run = src_cell.text_frame.paragraphs[0].runs[0]
                    if dest_cell.text_frame.paragraphs and dest_cell.text_frame.paragraphs[0].runs:
                        dest_run = dest_cell.text_frame.paragraphs[0].runs[0]
                        dest_run.font.size = src_run.font.size
                        dest_run.font.bold = src_run.font.bold
                        dest_run.font.name = src_run.font.name
                        dest_run.font.italic = src_run.font.italic
                        # Only copy RGB color if available, otherwise skip to avoid AttributeError
                        if src_run.font.color and src_run.font.color.type is not None:
                            try:
                                rgb = src_run.font.color.rgb
                                if rgb is not None:
                                    dest_run.font.color.rgb = rgb
                            except AttributeError:
                                # src_run.font.color may be a scheme color or other type without .rgb
                                pass
                # Copy fill color
                try:
                    fill = src_cell.fill
                    new_fill = dest_cell.fill
                    if fill.type is not None and fill.fore_color.rgb is not None:
                        new_fill.solid()
                        new_fill.fore_color.rgb = fill.fore_color.rgb
                    else:
                        new_fill.solid()
                        new_fill.fore_color.rgb = RGBColor(255, 255, 255)
                    dest_cell = apply_border(dest_cell, edges=["left", "right", "top", "bottom"], border_color="000000", border_width=2)
                except Exception:
                    pass

            # Copy the rest of the table structure (clear text for data rows)
            for r in range(1, rows):
                new_table.rows[r].height = table.rows[1].height
                for c in range(cols):
                    new_table.columns[c].width = table.columns[c].width
                    src_cell = table.cell(1, c)
                    dest_cell = new_table.cell(r, c)
                    dest_cell.text = ""  # Clear text
                    
                    # Copy fill color
                    try:
                        fill = src_cell.fill
                        new_fill = dest_cell.fill
                        if fill.type is not None and fill.fore_color.rgb is not None:
                            new_fill.solid()
                            new_fill.fore_color.rgb = fill.fore_color.rgb
                        else:
                            new_fill.solid()
                            new_fill.fore_color.rgb = RGBColor(255, 255, 255)
                        dest_cell = apply_border(dest_cell, edges=["left", "right", "top", "bottom"], border_color="000000", border_width=2)
                        dest_cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    except Exception:
                        pass
            break  # Only one table per slide is handled

    # Add all other shapes (except the table) from the source slide to the new slide
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in source_slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
            # Copy text boxes
            if shape.has_text_frame:
                textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                textbox_tf = textbox.text_frame
                # Remove all paragraphs so the text frame is truly empty
                while len(textbox_tf.paragraphs) > 0:
                    p = textbox_tf.paragraphs[0]
                    textbox_tf._element.remove(p._element)
                # Now add paragraphs from the source, skipping empty ones
                for p in shape.text_frame.paragraphs:
                    if not p.text.strip():
                        continue
                    new_p = textbox_tf.add_paragraph()
                    new_p.text = p.text
                    new_p.alignment = p.alignment
                    for i, run in enumerate(p.runs):
                        if i == 0:
                            dest_run = new_p.runs[0]
                        else:
                            dest_run = new_p.add_run()
                        dest_run.text = run.text
                        if run.font:
                            dest_font = dest_run.font
                            dest_font.name = run.font.name
                            dest_font.size = run.font.size
                            dest_font.bold = run.font.bold
                            dest_font.italic = run.font.italic
                            if run.font.color and run.font.color.type is not None:
                                dest_font.color.rgb = run.font.color.rgb
            # Copy pictures
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                from io import BytesIO
                image_stream = BytesIO(image_bytes)
                new_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)
            # Copy auto shapes (basic shapes)
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
                # Optionally copy fill color
                try:
                    if shape.fill.type is not None and shape.fill.fore_color.rgb is not None:
                        new_shape.fill.solid()
                        new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
                except Exception:
                    pass
    return prs

def load_editable_presentation(source_path: str, month_index: int) -> Presentation:
    """
    Copies the reference PPTX, opens it as a Presentation object, and replaces
    all instances of 'Month' with the specified month (0=January, 11=December).

    Returns the editable Presentation object so more slides can be added.
    """
    months = ["January", "February", "March", "April", "May", "June",
              "July", "August", "September", "October", "November", "December"]

    # Open the copied presentation
    prs = Presentation(source_path)

    # Replace 'Month' in all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace("Month", months[month_index])

    return prs  # You can now add slides, then save later


def generate_pbj_presentation(prs:Presentation,metadata:dict, month="May",report_type="Both"):
    
    for idx, (corp, corp_data) in enumerate(metadata.items()):
        corp_base = corp.lower().strip().replace(' ', '_')
        logo_dir = os.path.join(BASE_DIR, 'static', 'logos')
        logo_path = None
        for ext in [".png", ".jpg", ".jpeg"]:
            candidate = os.path.join(logo_dir, f"{corp_base}_logo{ext}")
            if os.path.exists(candidate):
                logo_path = candidate
                break
        if report_type == "overview" or report_type == "Both":
            # Duplicate the second slide as a template for corporate level
            df = corp_data['Corp_data']
            # Split the corporate summary table into slides of 10 facilities each
            facilities_per_slide = 10
            num_facilities = len(df)
            for start_idx in range(0, num_facilities, facilities_per_slide):
                end_idx = min(start_idx + facilities_per_slide, num_facilities)
                rows = end_idx - start_idx + 1  # +1 for header
                prs = duplicate_slide(prs, 2, rows)
                corp_slide = prs.slides[-1]
                for shape in corp_slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.text = run.text.replace("Corporate Name", corp)
                                run.text = run.text.replace("Month", month)
                    if shape.shape_type == 13:  # PICTURE
                        if shape.name.lower() == "picture 4":
                            if logo_path:
                                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                                sp = shape._element
                                corp_slide.shapes._spTree.remove(sp)
                                corp_slide.shapes.add_picture(logo_path, left, top, width, height)
                            else:
                                # Remove the logo shape from corp_slide if not found
                                sp = shape._element
                                corp_slide.shapes._spTree.remove(sp)
                                print(f"Logo not found for {corp}: {corp_base}_logo(.png/.jpg/.jpeg)")
                    if shape.has_table:
                        table = shape.table
                        fontSize = table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size
                        from pptx.enum.text import MSO_VERTICAL_ANCHOR
                        for idx, row in enumerate(df.iloc[start_idx:end_idx].itertuples()):
                            facilty = row.Facility
                            clinicians = row.Clinician
                            hours = row.Hours
                            hours = round(hours, 2) if hours is not None else 0.0
                            # Assuming the table has 3 columns: Facility, Clinicians, Hours
                            row_cells = table.rows[idx+1].cells
                            row_cells[0].text = facilty
                            row_cells[1].text = clinicians
                            row_cells[2].text = str(hours)
                            # Set vertical alignment to middle for all cells and all paragraphs in each cell
                            for cell in row_cells:
                                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        run.font.size = fontSize
        if report_type == "detailed" or report_type == "Both":
            # --- SPLIT FACILITIES ACROSS SLIDES BASED ON MAX PROVIDERS ---
            facilities = list(corp_data['Facility_Data'].items())
            max_providers_per_slide = 12
            def get_provider_count(facility_data):
                providers = facility_data['Provider']
                if isinstance(providers, (list, tuple, set)):
                    return len(providers)
                elif hasattr(providers, 'tolist'):
                    return len(providers.tolist())
                elif isinstance(providers, str):
                    return 1 if providers.strip() else 0
                else:
                    return 1

            # Calculate total providers
            total_providers = sum(get_provider_count(facility_data) for _, facility_data in facilities)
            facilities_per_slide = 5 if total_providers < max_providers_per_slide else None

            slide_batches = []
            current_slide_facilities = []
            current_provider_count = 0
            i = 0
            while i < len(facilities):
                facility, facility_data = facilities[i]
                count = get_provider_count(facility_data)
                # If a facility has more than max_providers_per_slide, try to add another facility with it
                if count > max_providers_per_slide-5:
                    # Try to add the next facility if possible
                    batch = [(facility, facility_data)]
                    if i + 1 < len(facilities):
                        next_facility, next_facility_data = facilities[i+1]
                        next_count = get_provider_count(next_facility_data)
                        if count + next_count <= max_providers_per_slide+10:
                            batch.append((next_facility, next_facility_data))
                            i += 1
                    slide_batches.append(batch)
                    i += 1
                    continue
                # If limiting by facilities_per_slide (when total providers is small)
                if facilities_per_slide:
                    current_slide_facilities.append((facility, facility_data))
                    if len(current_slide_facilities) == facilities_per_slide:
                        slide_batches.append(current_slide_facilities)
                        current_slide_facilities = []
                    i += 1
                    continue
                # Otherwise, use provider count batching
                if current_provider_count + count > max_providers_per_slide and current_slide_facilities:
                    slide_batches.append(current_slide_facilities)
                    current_slide_facilities = []
                    current_provider_count = 0
                current_slide_facilities.append((facility, facility_data))
                current_provider_count += count
                i += 1
            if current_slide_facilities:
                slide_batches.append(current_slide_facilities)

            for batch_idx, batch in enumerate(slide_batches):
                prs = duplicate_slide(prs, 3, len(batch) + 1)
                facility_slide = prs.slides[-1]
                for shape in facility_slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.text = run.text.replace("Corporate Name", corp)
                                run.text = run.text.replace("Month", month)
                    if shape.shape_type == 13:
                        if shape.name.lower() == "picture 5":
                            # Use the same logo_path logic as above
                            if logo_path:
                                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                                sp = shape._element
                                facility_slide.shapes._spTree.remove(sp)
                                facility_slide.shapes.add_picture(logo_path, left, top, width, height)
                            else:
                                # Remove the logo shape from all facility slides if not found
                                sp = shape._element
                                facility_slide.shapes._spTree.remove(sp)
                                print(f"Logo not found for {corp}: {corp_base}_logo(.png/.jpg/.jpeg)")
                    if shape.has_table:
                        table = shape.table
                        fontSize=table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size
                        font=table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.name
                        from pptx.enum.text import MSO_VERTICAL_ANCHOR
                        for idx, (facility, facility_data) in enumerate(batch):
                            # Ensure the table has enough rows before accessing
                            if idx + 1 >= len(table.rows):
                                # Not enough rows, skip or break to avoid IndexError
                                print(f"Warning: Not enough table rows for facility data at index {idx+1}. Skipping.")
                                break
                            # --- original facility row logic ---
                            providers = facility_data['Provider']
                            if isinstance(providers, (list, tuple, set)):
                                providers_str = '\n'.join(str(p) for p in providers)
                            elif hasattr(providers, 'tolist'):
                                providers_str = '\n'.join(str(p) for p in providers.tolist())
                            else:
                                providers_str = str(providers)
                            dos = facility_data['DOS']
                            if isinstance(dos, (list, tuple, set)):
                                dos_str = '\n'.join(str(d) for d in dos)
                            elif hasattr(dos, 'tolist'):
                                dos_str = '\n'.join(str(d) for d in dos.tolist())
                            else:
                                dos_str = str(dos)
                            hours = facility_data['Hours']
                            try:
                                hours_val = round(float(hours), 2)
                            except Exception:
                                hours_val = hours
                            total_hours = None
                            if hasattr(hours, 'sum'):
                                try:
                                    total_hours = round(float(hours.sum()), 2)
                                except Exception:
                                    total_hours = hours
                            elif isinstance(hours, (list, tuple, set)):
                                try:
                                    total_hours = round(sum(float(h) for h in hours), 2)
                                except Exception:
                                    total_hours = hours
                            else:
                                total_hours = hours_val
                            if isinstance(hours, (list, tuple, set)) or hasattr(hours, 'tolist'):
                                if hasattr(hours, 'tolist'):
                                    hours_list = [str(round(float(h), 2)) for h in hours.tolist()]
                                else:
                                    hours_list = [str(round(float(h), 2)) for h in hours]
                                hours_str = '\n'.join(hours_list)
                            else:
                                hours_str = str(hours_val)
                            row_cells = table.rows[idx+1].cells
                            row_cells[0].text = str(facility)
                            providers_lines = providers_str.split('\n') if providers_str else []
                            row_cells[1].text = '\n'.join(providers_lines)
                            empty_paragraph = row_cells[1].text_frame.add_paragraph()
                            empty_paragraph.text = '\n.'  # Add an empty paragraph to ensure proper spacing
                            row_cells[3].text = hours_str
                            para_total_val = row_cells[3].text_frame.add_paragraph()
                            para_total_val.text = f"\n{str(total_hours)}"
                            para_total_val.runs[0].font.bold = True
                            row_cells[2].text = dos_str
                            para_total = row_cells[2].text_frame.add_paragraph()
                            para_total.text = '\nTotal:'
                            para_total.runs[0].font.bold = True
                            for cell in row_cells:
                                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        run.font.size = fontSize
                                        run.font.name = font
    remove_slide(prs, prs.slides[1])  # Remove the first slide if it's a title slide or not needed
    remove_slide(prs, prs.slides[1])  # Remove the second slide (after first removal, next is at index 1)


@app.route('/')
def index():
    return render_template('PPG_report.html')

# Flask API endpoint to receive file and CPT dict
@app.route('/upload_data', methods=['POST'])
def upload_data():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    file = request.files['file']
    filename = secure_filename(file.filename)
    ext = os.path.splitext(filename)[1].lower()
    if ext not in ['.csv', '.xlsx']:
        return jsonify({'error': 'Invalid file type'}), 400
    # Save file to a temp location
    temp_path = os.path.join('uploads', filename)
    os.makedirs('uploads', exist_ok=True)
    file.save(temp_path)
    # Load DataFrame
    if ext == '.csv':
        charge_capture_df = pd.read_csv(temp_path)
    else:
        charge_capture_df = pd.read_excel(temp_path)
    # Get CPT dict
    cpt_time_dict = request.form.get('cpt_time_dict')
    import json
    if cpt_time_dict:
        cpt_time_dict = json.loads(cpt_time_dict)
    else:
        cpt_time_dict = {
            "99304": 60,
            "99305": 75,
            "99306": 90,
            "99307": 15,
            "99308": 25,
            "99309": 35,
            "99310": 45
        }

    # Get month and report_type from form
    month_index = request.form.get('month')
    try:
        month_index = int(month_index)
    except (TypeError, ValueError):
        month_index = 4  # Default to May if not provided
    months = ["January", "February", "March", "April", "May", "June",
              "July", "August", "September", "October", "November", "December"]
    month_name = months[month_index] if 0 <= month_index < 12 else "May"
    report_type = request.form.get('report_type', 'Both')

    # Process as before
    corporate_groups, single_facilities = group_facilities(charge_capture_df)
    metadata = build_metadata(charge_capture_df, corporate_groups, single_facilities)
    prs = load_editable_presentation(os.path.join(BASE_DIR, 'static', 'reference_slide.pptx'), month_index=month_index)
    generate_pbj_presentation(prs, metadata, month=month_name, report_type=report_type)
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], 'generated_pbj_report.pptx')
    prs.save(output_path)
    # Optionally, return download link or status
    # Return the relative path for the frontend to use in download URL
    return jsonify({'success': True, 'pptx_path': f"uploads/generated_pbj_report.pptx", 'month': month_name, 'report_type': report_type})


from flask import send_from_directory

# Serve files from the uploads directory
@app.route('/uploads/<path:filename>')
def download_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename, as_attachment=True)

# API endpoint to upload a logo image for a corporate group
@app.route('/upload_logo', methods=['POST'])
def upload_logo():
    if 'logo' not in request.files or 'corp_name' not in request.form:
        return jsonify({'error': 'Logo file and corp_name are required'}), 400
    logo_file = request.files['logo']
    corp_name = request.form['corp_name'].strip().lower().replace(' ', '_')
    ext = os.path.splitext(logo_file.filename)[1].lower()
    if ext not in ['.jpg', '.jpeg', '.png']:
        return jsonify({'error': 'Invalid file type'}), 400
    logos_dir = os.path.join('static', 'logos')
    os.makedirs(logos_dir, exist_ok=True)
    logo_filename = f"{corp_name}_logo{ext}"
    logo_path = os.path.join(logos_dir, logo_filename)
    logo_file.save(logo_path)
    return jsonify({'success': True, 'logo_path': logo_path})

# Function to open the default web browser
def open_browser():
    webbrowser.open_new("http://localhost:9000")

if __name__ == '__main__':
    threading.Timer(1, open_browser).start()  # Open the browser after 1 second
    app.run(host='localhost', port=9000)