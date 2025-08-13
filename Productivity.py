import pandas as pd
import os
from flask import Flask, request, jsonify, render_template
from werkzeug.utils import secure_filename
import sys
import threading
import webbrowser
import re

from rapidfuzz import process, fuzz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import  Pt
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement

def get_base_dir():
    if getattr(sys, 'frozen', False):
        return sys._MEIPASS  # PyInstaller extracts to this temp dir
    return os.path.dirname(os.path.abspath(__file__))

app = Flask(__name__,
            template_folder=os.path.join(get_base_dir(), 'templates'))

BASE_DIR = get_base_dir()
print("Base Directory:", BASE_DIR)
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)


def remove_slide(prs, slide):
    # Helper to remove a slide from a Presentation (python-pptx workaround)
    slide_id = slide.slide_id
    slides = prs.slides._sldIdLst
    for sldId in slides:
        if int(sldId.attrib['id']) == slide_id:
            slides.remove(sldId)
            break


def match_providers(company_roster_providers, charge_capture_providers, threshold=76):
    matched = []
    unmatched = []
    already_matched_charge_capture_names = set()

    # Create mappings for normalized to original names
    company_roster_map = {str(provider).strip().lower(): provider for provider in company_roster_providers}
    charge_capture_map = {str(provider).strip().lower(): provider for provider in charge_capture_providers}

    # Normalize the lists
    company_roster_providers = list(company_roster_map.keys())
    charge_capture_providers = list(charge_capture_map.keys())

    for current_threshold in range(92, threshold - 1, -4):  # Iterate thresholds inclusively
        for company_roster in company_roster_providers:
            # Skip if already matched
            if company_roster_map[company_roster] in [item['company_roster'] for item in matched]:
                continue

            # Filter out already matched charge capture names
            available_charge_capture_providers = [
                name for name in charge_capture_providers if name not in already_matched_charge_capture_names
            ]
            match_name, score, _ = process.extractOne(
                company_roster,
                available_charge_capture_providers,
                scorer=fuzz.ratio
            )

            if score >= current_threshold:
                matched.append({
                    'company_roster': company_roster_map[company_roster],  # Use original name
                    'charge_capture_name': charge_capture_map[match_name],  # Use original name
                    'score': score
                })
                already_matched_charge_capture_names.add(match_name)
            else:
                if company_roster_map[company_roster] not in [item['company_roster'] for item in unmatched]:
                    unmatched.append({
                        'company_roster': company_roster_map[company_roster],  # Use original name
                        'charge_capture_name': charge_capture_map[match_name],  # Use original name
                        'score': score
                    })
    filtered_unmatched = [item for item in unmatched if item['company_roster'] not in [m['company_roster'] for m in matched]]
    unmatched = filtered_unmatched
    return {"matched": matched, "unmatched": unmatched}

def apply_border(cell, edges = ["left", "right", "top", "bottom"], border_color="000000", border_width=1):
    if type(edges) is not list: edges = [edges]
    border_width = str(border_width*Pt(1))
    def SubElement(parent, tagname, **kwargs):
            element = OxmlElement(tagname)
            element.attrib.update(kwargs)
            parent.append(element)
            return element
    
    lines = [{"left": 'a:lnL',
              "right": 'a:lnR',
              "top": 'a:lnT',
              "bottom": 'a:lnB'}[_] for _ in edges]
    
    if cell.fill.type == MSO_FILL.SOLID: fill_color = cell.fill.fore_color.rgb
    cell.fill.background()
    
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for line in lines:
        
        # Remove duplicate tag if it exists
        tag = line.split(":")[-1]
        for e in tcPr.getchildren():
            if tag in str(e.tag): tcPr.remove(e)
        
        ln = SubElement(tcPr, line, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
        
    cell.fill.solid()
    if cell.fill.type == MSO_FILL.SOLID: cell.fill.fore_color.rgb = fill_color
    return(cell)

def duplicate_slide(prs: Presentation, slide_number: int,rows: int):
    """
    Duplicate a slide, copying only tables (structure and cell colors, no text).
    :param prs: Presentation object.
    :param slide_number: 1-based index of the slide to duplicate.
    :return: The newly created slide.
    """
    index = slide_number - 1
    if index < 0 or index >= len(prs.slides):
        raise IndexError(f"Slide {slide_number} does not exist.")

    source_slide = prs.slides[index]
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Remove all shapes from the new slide
    for shape in list(new_slide.shapes):
        sp = shape._element
        new_slide.shapes._spTree.remove(sp)

    # Add the table from the source slide
    for shape in source_slide.shapes:
        if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
            table = shape.table
            rows, cols = rows, len(table.columns)
            new_table_shape = new_slide.shapes.add_table(rows, cols, shape.left, shape.top, shape.width, shape.height)
            new_table = new_table_shape.table
            # Ensure header row height matches the source header row height
            new_table.rows[0].height = table.rows[0].height

            # Copy the first row (header) exactly
            from pptx.enum.text import MSO_VERTICAL_ANCHOR
            for c in range(cols):
                src_cell = table.cell(0, c)
                dest_cell = new_table.cell(0, c)
                dest_cell.text = src_cell.text
                # Set vertical alignment to middle
                dest_cell.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                # Copy font properties from the first run of the first paragraph if available
                if src_cell.text_frame.paragraphs and src_cell.text_frame.paragraphs[0].runs:
                    src_run = src_cell.text_frame.paragraphs[0].runs[0]
                    if dest_cell.text_frame.paragraphs and dest_cell.text_frame.paragraphs[0].runs:
                        dest_run = dest_cell.text_frame.paragraphs[0].runs[0]
                        dest_run.font.size = src_run.font.size
                        dest_run.font.bold = src_run.font.bold
                        dest_run.font.name = src_run.font.name
                        dest_run.font.italic = src_run.font.italic
                        # Only copy RGB color if available, otherwise skip to avoid AttributeError
                        if src_run.font.color and src_run.font.color.type is not None:
                            try:
                                rgb = src_run.font.color.rgb
                                if rgb is not None:
                                    dest_run.font.color.rgb = rgb
                            except AttributeError:
                                # src_run.font.color may be a scheme color or other type without .rgb
                                pass
                # Copy fill color
                try:
                    fill = src_cell.fill
                    new_fill = dest_cell.fill
                    if fill.type is not None and fill.fore_color.rgb is not None:
                        new_fill.solid()
                        new_fill.fore_color.rgb = fill.fore_color.rgb
                    else:
                        new_fill.solid()
                        new_fill.fore_color.rgb = RGBColor(255, 255, 255)
                    dest_cell = apply_border(dest_cell, edges=["left", "right", "top", "bottom"], border_color="000000", border_width=2)
                except Exception:
                    pass

            # Copy the rest of the table structure (clear text for data rows)
            for r in range(1, rows):
                new_table.rows[r].height = table.rows[1].height
                for c in range(cols):
                    new_table.columns[c].width = table.columns[c].width
                    src_cell = table.cell(1, c)
                    dest_cell = new_table.cell(r, c)
                    dest_cell.text = ""  # Clear text
                    
                    # Copy fill color
                    try:
                        fill = src_cell.fill
                        new_fill = dest_cell.fill
                        if fill.type is not None and fill.fore_color.rgb is not None:
                            new_fill.solid()
                            new_fill.fore_color.rgb = fill.fore_color.rgb
                        else:
                            new_fill.solid()
                            new_fill.fore_color.rgb = RGBColor(255, 255, 255)
                        dest_cell = apply_border(dest_cell, edges=["left", "right", "top", "bottom"], border_color="000000", border_width=2)
                        dest_cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    except Exception:
                        pass
            break  # Only one table per slide is handled

    # Add all other shapes (except the table) from the source slide to the new slide
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in source_slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
            # Copy text boxes
            if shape.has_text_frame:
                textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                textbox_tf = textbox.text_frame
                # Remove all paragraphs so the text frame is truly empty
                while len(textbox_tf.paragraphs) > 0:
                    p = textbox_tf.paragraphs[0]
                    textbox_tf._element.remove(p._element)
                # Now add paragraphs from the source, skipping empty ones
                for p in shape.text_frame.paragraphs:
                    if not p.text.strip():
                        continue
                    new_p = textbox_tf.add_paragraph()
                    new_p.text = p.text
                    new_p.alignment = p.alignment
                    for i, run in enumerate(p.runs):
                        if i == 0:
                            dest_run = new_p.runs[0]
                        else:
                            dest_run = new_p.add_run()
                        dest_run.text = run.text
                        if run.font:
                            dest_font = dest_run.font
                            dest_font.name = run.font.name
                            dest_font.size = run.font.size
                            dest_font.bold = run.font.bold
                            dest_font.italic = run.font.italic
                            if run.font.color and run.font.color.type is not None:
                                dest_font.color.rgb = run.font.color.rgb
            # Copy pictures
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                from io import BytesIO
                image_stream = BytesIO(image_bytes)
                new_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)
            # Copy auto shapes (basic shapes)
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
                # Optionally copy fill color
                try:
                    if shape.fill.type is not None and shape.fill.fore_color.rgb is not None:
                        new_shape.fill.solid()
                        new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
                except Exception:
                    pass
    return prs

def load_editable_presentation(source_path: str, day: str , Date: str) -> Presentation:
    """
    Copies the reference PPTX, opens it as a Presentation object, and replaces
    all instances of 'Month' with the specified month (0=January, 11=December).

    Returns the editable Presentation object so more slides can be added.
    """

    # Open the copied presentation
    prs = Presentation(source_path)

    # Replace 'Month' in all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace("DAY", day)
                        run.text = run.text.replace("Date", Date)

    return prs  # You can now add slides, then save later

def build_metadata(charge_capture_df, company_roster_df):
    CC_used_col = ["Provider", "CPT Codes", "Charge Status"]
    CR_used_col = ['Name', "Manager", 'State/Region']

    CC_filtered_df = charge_capture_df[CC_used_col]
    CR_filtered_df = company_roster_df[CR_used_col]

    # Rename column safely
    CR_filtered_df = CR_filtered_df.rename(columns={'Name': 'Provider'})

    grouped_CC = CC_filtered_df.groupby('Provider')
    grouped_CR = CR_filtered_df.groupby('State/Region')
    cpt_pattern = re.compile(r'\b993\d{2}\b')  # Example pattern for 6-digit codes starting with 9930
    Region_data = pd.DataFrame(columns=['Clinician', 'Gross Encounters', 'C99304', 'C99305', 'C99306', 'C99307', 'C99308', 'C99309', 'C99310','Drafts','CCM_counts',"manager","region"])
    for name, group in grouped_CC:
    
        group['CPT Codes'] = group['CPT Codes'].astype(str).str.split(',')
        group = group.explode('CPT Codes')
        group['CPT Codes'] = group['CPT Codes'].str.strip()

        group['Charge Status'] = group['Charge Status'].str.strip().str.lower()
        # CCM Counts
        CCM_counts =group[group['CPT Codes']=="44444"]
        CCM_counts = CCM_counts['CPT Codes'].count()

        # Filter rows where 'CPT Codes' is in the target list and 'Charge Status' is 'draft'
        draft_counts = group[
            (group['CPT Codes'].str.match(cpt_pattern, na=False)) & 
            (group['Charge Status'] == 'draft')
        ]

        # Count the rows
        draft_counts = draft_counts['Charge Status'].count()
        # Gross Encounters    
        cpt_grouped = group[group['CPT Codes'].str.match(cpt_pattern, na=False)]
        gross_encounters = cpt_grouped['CPT Codes'].count()

        new_row = {'Clinician': name, 'Gross Encounters': gross_encounters,
                   'C99304': cpt_grouped['CPT Codes'].str.contains('99304').sum(),
                   'C99305': cpt_grouped['CPT Codes'].str.contains('99305').sum(),
                    'C99306': cpt_grouped['CPT Codes'].str.contains('99306').sum(),
                    'C99307': cpt_grouped['CPT Codes'].str.contains('99307').sum(),
                    'C99308': cpt_grouped['CPT Codes'].str.contains('99308').sum(),
                    'C99309': cpt_grouped['CPT Codes'].str.contains('99309').sum(),
                    'C99310': cpt_grouped['CPT Codes'].str.contains('99310').sum(),
                    'Drafts': draft_counts,
                    'CCM_counts': CCM_counts}
        Region_data.loc[len(Region_data)] = new_row
    unmatched_providers = []
    matched_providers = []
    Regional_Dashboard=pd.DataFrame(columns=["Region","RDO","RDS","Gross Encounters","Gross Consents","Gross Drafts"])

    for name, group in grouped_CR:
        manager_group=group.groupby('Manager')
        for manager_name, manager_group in manager_group:
            Clinician_list=Region_data['Clinician']
            Clinician_list2 = manager_group['Provider']
            result=match_providers(Clinician_list2, Clinician_list, threshold=74)
            unmatched_providers.extend(result['unmatched'])
            matched_providers.extend(result['matched'])
            Reagion_Gross_Encounters = Region_data[Region_data['Clinician'].isin([item['charge_capture_name'] for item in result['matched']])]['Gross Encounters'].sum()
            Reagion_Gross_Consents = Region_data[Region_data['Clinician'].isin([item['charge_capture_name'] for item in result['matched']])]['CCM_counts'].sum()
            Reagion_Gross_Drafts = Region_data[Region_data['Clinician'].isin([item['charge_capture_name'] for item in result['matched']])]['Drafts'].sum()    
            new_row = {
                "Region": name,
                "RDO": manager_name,
                "RDS": Clinician_list2.to_list(),
                "Gross Encounters": Reagion_Gross_Encounters,
                "Gross Consents": Reagion_Gross_Consents,
                "Gross Drafts": Reagion_Gross_Drafts
            }
            Regional_Dashboard.loc[len(Regional_Dashboard)] = new_row
    Regional_Dashboard.columns = Regional_Dashboard.columns.str.replace(' ', '_')
    Region_data.columns = Region_data.columns.str.replace(' ', '_')
    for row in Regional_Dashboard.itertuples():
        region = row.Region
        Clinician_list = row.RDS
        manager= row.RDO    
        # Update the 'region' column for all clinicians in the list
        Region_data.loc[Region_data['Clinician'].isin(Clinician_list), 'region'] = region
        Region_data.loc[Region_data['Clinician'].isin(Clinician_list), 'manager'] = manager

    
    # Drop rows with NaN values in the 'region' column
    Region_data.dropna(subset=['region'], inplace=True)
    
    gouped_CC_ByRegion = Region_data.groupby(['region', 'manager'])
    return Regional_Dashboard, unmatched_providers, matched_providers, gouped_CC_ByRegion
def generate_pbj_presentation(prs:Presentation,Regional_Dashboard: pd.DataFrame,gouped_CC_ByRegion, day:str, date:str,Add_region=True):
    
    # Split the corporate summary table into slides of 10 facilities each
    Region_per_slide = 5
    num_Region = len(Regional_Dashboard)
    for start_idx in range(0, num_Region, Region_per_slide):
        end_idx = min(start_idx + Region_per_slide, num_Region)
        rows = end_idx - start_idx + 1  # +1 for header
        prs = duplicate_slide(prs, 2, rows)
        region_dashborad_slide = prs.slides[-1]
        for shape in region_dashborad_slide.shapes:
            if shape.has_table:
                table = shape.table
                fontSize = table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size
                from pptx.enum.text import MSO_VERTICAL_ANCHOR
                for idx, row in enumerate(Regional_Dashboard.iloc[start_idx:end_idx].itertuples()):
                    region = row.Region
                    RDO = row.RDO
                    RDS = ', '.join(map(str, row.RDS))  # Convert list to comma-separated string
                    gross_encounters = row.Gross_Encounters
                    gross_consents = row.Gross_Consents
                    gross_drafts = row.Gross_Drafts
                    # Assuming the table has 3 columns: Facility, Clinicians, Hours
                    row_cells = table.rows[idx+1].cells
                    row_cells[0].text = str(region)
                    row_cells[1].text = str(RDO)
                    # Split RDS into two columns, pad each provider name to fixed width, and set monospace font
                    monospace_font = "Consolas"
                    pad_width = 25
                    names = [str(name) if len(str(name)) > pad_width else str(name).ljust(pad_width) for name in row.RDS]
                    col1 = []
                    col2 = []
                    for i, name in enumerate(names):
                        if i % 2 == 0:
                            col1.append(name)
                        else:
                            col2.append(name)
                    # Make both columns the same length
                    max_len = max(len(col1), len(col2))
                    while len(col1) < max_len:
                        col1.append(' ' * pad_width)
                    while len(col2) < max_len:
                        col2.append('')
                    lines = [col1[i] + col2[i] for i in range(max_len)]
                    combined_text = '\n'.join(lines)
                    row_cells[2].text = combined_text
                    # Set monospace font for all runs in the cell
                    for para in row_cells[2].text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = monospace_font
                    row_cells[3].text = str(gross_encounters)
                    row_cells[4].text = str(gross_consents)
                    row_cells[6].text = str(gross_drafts)
                    # Set vertical alignment to middle for all cells and all paragraphs in each cell
                    for cell in row_cells:
                        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.size = fontSize
    for region_name , region_data in gouped_CC_ByRegion:
        # --- SPLIT PROVIDERS ACROSS SLIDES BASED ON MAX PROVIDERS ---
        max_providers_per_slide = 7
        providers_df = region_data.reset_index(drop=True)
        total_providers = len(providers_df)
        print(f"Region: {region_name}, Total Providers: {total_providers}")

        # Create as many slides as needed to fit all providers
        for start in range(0, total_providers, max_providers_per_slide):
            end = min(start + max_providers_per_slide, total_providers)
            rows = (end - start) + 1  # +1 for header
            prs = duplicate_slide(prs, 3, rows)
            Region_Scope_slide = prs.slides[-1]

            for shape in Region_Scope_slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = run.text.replace("Region", f"{region_name[0]} - Managers: {region_name[1]}")
                if shape.has_table:
                    table = shape.table
                    fontSize = table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size
                    font = table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.name
                    from pptx.enum.text import MSO_VERTICAL_ANCHOR

                    # Fill only the current chunk of providers
                    for idx, row in enumerate(providers_df.iloc[start:end].itertuples(index=False)):
                        clinician = row.Clinician
                        gross_encounters = row.Gross_Encounters
                        C993_04 = row.C99304
                        C993_05 = row.C99305
                        C993_06 = row.C99306
                        C993_07 = row.C99307
                        C993_08 = row.C99308
                        C993_09 = row.C99309
                        C993_10 = row.C99310
                        gross_consents = row.CCM_counts
                        gross_drafts = row.Drafts

                        row_cells = table.rows[idx + 1].cells
                        row_cells[0].text = str(clinician)
                        row_cells[1].text = str(gross_encounters)
                        row_cells[2].text = str(C993_04)
                        row_cells[3].text = str(C993_05)
                        row_cells[4].text = str(C993_06)
                        row_cells[5].text = str(C993_07)
                        row_cells[6].text = str(C993_08)
                        row_cells[7].text = str(C993_09)
                        row_cells[8].text = str(C993_10)
                        row_cells[10].text = str(gross_drafts)
                        row_cells[11].text = str(gross_consents)

                        # Set vertical alignment to middle for all cells and font size
                        for cell in row_cells:
                            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    run.font.size = fontSize
    remove_slide(prs, prs.slides[1])  # Remove the first slide if it's a title slide or not needed
    remove_slide(prs, prs.slides[1])  # Remove the second slide (after first removal, next is at index 1)

@app.route('/')
def index():
    return render_template('Weekly_report.html')

# Flask API endpoint to receive file and CPT dict
@app.route('/upload_data', methods=['POST'])
def upload_data():
    if 'Roster' not in request.files or 'captureFile' not in request.files:
        return jsonify({'error': 'Missing files'}), 400
    Roster_file = request.files['Roster']
    capture_file = request.files['captureFile']

    capture_file_filename = secure_filename(capture_file.filename)
    Roster_file_filename = secure_filename(Roster_file.filename)
    captuer_ext = os.path.splitext(capture_file_filename)[1].lower()
    roster_ext = os.path.splitext(Roster_file_filename)[1].lower()
    if captuer_ext not in ['.csv', '.xlsx'] or roster_ext not in ['.csv', '.xlsx']:
        return jsonify({'error': 'Invalid file type'}), 400

    # Save files to a temp location
    # Save file to a temp location
    capture_temp_path = os.path.join(app.config['UPLOAD_FOLDER'], capture_file_filename)
    roster_temp_path = os.path.join(app.config['UPLOAD_FOLDER'], Roster_file_filename)

    capture_file.save(capture_temp_path)
    Roster_file.save(roster_temp_path)
    # Load DataFrame

    charge_capture_df = pd.read_excel(capture_temp_path)

    conpany_roaster = pd.read_excel(roster_temp_path,skiprows=2)


    # Get month and report_type from form
    day = request.form.get('day')
    date= request.form.get('date')
   
    include_regions = request.form.get('include_regions',True)

    Regional_Dashboard, unmatched_providers, matched_providers, gouped_CC_ByRegion = build_metadata(charge_capture_df, conpany_roaster)
    prs = load_editable_presentation(os.path.join(BASE_DIR, 'static', 'Weekly Report Example.pptx'), day= day, Date=date)
    generate_pbj_presentation(prs, Regional_Dashboard, gouped_CC_ByRegion, day, date, Add_region=include_regions)
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], 'Weekly_report_report.pptx')
    prs.save(output_path)
    # Optionally, return download link or status
    # Return the relative path for the frontend to use in download URL
    return jsonify({
        'success': True,
        'pptx_path': f"uploads/Weekly_report_report.pptx",
        'matched_providers': matched_providers,
        'unmatched_providers': unmatched_providers
    })


from flask import send_from_directory

# Serve files from the uploads directory
@app.route('/uploads/<path:filename>')
def download_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename, as_attachment=True)

# Function to open the default web browser
import os
def open_browser():
    webbrowser.open_new("http://localhost:8000")

if __name__ == '__main__':
    # Only open browser in local development
    if os.environ.get('RENDER') is None:
        threading.Timer(1, open_browser).start()  # Open the browser after 1 second
    port = int(os.environ.get('PORT', 8000))
    app.run(host='0.0.0.0', port=port)